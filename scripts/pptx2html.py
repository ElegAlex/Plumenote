#!/usr/bin/env python3
"""Convert PPTX to structured HTML with extracted images.

Usage: pptx2html.py <input.pptx> <media_dir>
Outputs HTML to stdout. Images are saved to <media_dir>/media/.
"""
import sys
import os
import base64
import hashlib
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print("ERROR: python-pptx not installed", file=sys.stderr)
    sys.exit(1)


def shape_text(shape):
    """Extract text from a shape's text_frame, preserving paragraphs."""
    if not shape.has_text_frame:
        return ""
    parts = []
    for para in shape.text_frame.paragraphs:
        runs_html = []
        for run in para.runs:
            text = run.text
            if not text:
                continue
            # Apply formatting
            if run.font.bold:
                text = f"<strong>{text}</strong>"
            if run.font.italic:
                text = f"<em>{text}</em>"
            if run.font.underline:
                text = f"<u>{text}</u>"
            if run.hyperlink and run.hyperlink.address:
                text = f'<a href="{run.hyperlink.address}">{text}</a>'
            runs_html.append(text)
        line = "".join(runs_html).strip()
        if line:
            parts.append(line)
    return parts


def process_table(table):
    """Convert a table shape to HTML table."""
    html = "<table>\n"
    for i, row in enumerate(table.rows):
        html += "  <tr>\n"
        tag = "th" if i == 0 else "td"
        for cell in row.cells:
            text = cell.text.strip()
            html += f"    <{tag}>{text}</{tag}>\n"
        html += "  </tr>\n"
    html += "</table>\n"
    return html


def main():
    if len(sys.argv) < 3:
        print(f"Usage: {sys.argv[0]} <input.pptx> <media_dir>", file=sys.stderr)
        sys.exit(1)

    pptx_path = sys.argv[1]
    media_dir = os.path.join(sys.argv[2], "media")
    os.makedirs(media_dir, exist_ok=True)

    prs = Presentation(pptx_path)
    html_parts = []
    img_counter = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_title = None
        slide_content = []

        # Sort shapes by position (top to bottom, left to right)
        shapes = sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))

        for shape in shapes:
            # Title
            if shape.has_text_frame and shape.shape_type in (
                MSO_SHAPE_TYPE.PLACEHOLDER,
            ) or (hasattr(shape, "is_placeholder") and shape.is_placeholder):
                ph = shape.placeholder_format
                if ph and ph.idx == 0:  # Title placeholder
                    title_parts = shape_text(shape)
                    if title_parts:
                        slide_title = title_parts[0]
                    continue

            # Images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                ext = img.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                if ext in ("x-emf", "x-wmf", "emf", "wmf"):
                    ext = "png"  # skip non-web formats
                    continue
                img_counter += 1
                filename = f"slide{slide_num}_img{img_counter}.{ext}"
                img_path = os.path.join(media_dir, filename)
                with open(img_path, "wb") as f:
                    f.write(img.blob)
                slide_content.append(
                    f'<img src="{img_path}" alt="{shape.name}" />'
                )
                continue

            # Tables
            if shape.has_table:
                slide_content.append(process_table(shape.table))
                continue

            # Group shapes — extract text from children
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    if child.has_text_frame:
                        for line in shape_text(child):
                            slide_content.append(f"<p>{line}</p>")
                    if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        img = child.image
                        ext = img.content_type.split("/")[-1]
                        if ext == "jpeg":
                            ext = "jpg"
                        if ext in ("x-emf", "x-wmf", "emf", "wmf"):
                            continue
                        img_counter += 1
                        filename = f"slide{slide_num}_img{img_counter}.{ext}"
                        img_path = os.path.join(media_dir, filename)
                        with open(img_path, "wb") as f:
                            f.write(img.blob)
                        slide_content.append(
                            f'<img src="{img_path}" alt="{child.name}" />'
                        )
                continue

            # Regular text shapes
            if shape.has_text_frame:
                paragraphs = shape_text(shape)
                for line in paragraphs:
                    # Detect bullet-like patterns
                    slide_content.append(f"<p>{line}</p>")

        # Build slide HTML
        title = slide_title or f"Slide {slide_num}"
        html_parts.append(f"<h2>{title}</h2>")
        if slide_content:
            html_parts.extend(slide_content)
        else:
            html_parts.append("<p></p>")

    print("\n".join(html_parts))


if __name__ == "__main__":
    main()
